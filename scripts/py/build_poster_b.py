"""Assemble PLAN B: the conventional four-section (IMRaD) 48 x 36 in poster.

This is the safe, report-style alternative to the three-act narrative in
build_poster_a.py. Same verified content, same figures, same header and
footer -- but a formal title and the standard Introduction / Methods /
Results / Conclusions structure a reviewer expects.

Layout: four columns. Introduction, Methods, then Results spanning columns
3 and 4, with Conclusions as a full-width band across the bottom.

python-pptx cannot measure rendered text, so block heights are estimated from
character counts and each column's running y-position is tracked explicitly.
The script prints a per-column budget at the end -- if a column exceeds BODY_H
it will collide with the conclusions band, so check that output.

Run:  uv run --python 3.12 --with python-pptx --with pillow python \
          scripts/py/build_poster_b.py
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
FIG = ROOT / "figures" / "poster"
LOGO = ROOT / "BINF6999" / "drafts" / "logos"
OUT = ROOT / "BINF6999" / "drafts" / "LiF_BINF6999_Poster_Draft_B.pptx"

W, H = 48.0, 36.0
FONT = "Segoe UI"

INK = RGBColor(0x16, 0x20, 0x2A)
INK_D = RGBColor(0x0E, 0x16, 0x1E)
INK_2 = RGBColor(0x5C, 0x6B, 0x7A)
TEAL = RGBColor(0x0E, 0x7C, 0x6B)
RED = RGBColor(0xBE, 0x3A, 0x34)
SALMON = RGBColor(0xE8, 0x8A, 0x84)
PAPER = RGBColor(0xFC, 0xFC, 0xFA)
BAND = RGBColor(0xEF, 0xF2, 0xF1)
LINE = RGBColor(0xD8, 0xDE, 0xE2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9A, 0xAA, 0xB8)
TINT = RGBColor(0xF4, 0xF7, 0xF6)

MARGIN, GUT = 1.30, 0.76
COL_W = (W - 2 * MARGIN - 3 * GUT) / 4
COL_X = [MARGIN + i * (COL_W + GUT) for i in range(4)]

HDR_H = 4.85
BODY_Y = 5.45
BODY_H = 22.30
GAP = 0.24


# ---------------------------------------------------------------- text metrics
def est_h(txt: str, w: float, size: float, spacing: float = 1.15,
          bold: bool = False) -> float:
    char_w = 0.545 if bold else 0.495
    cpl = max(1, int(w * 72 / (char_w * size)))
    lines = sum(max(1, -(-len(seg) // cpl)) for seg in txt.split("\n"))
    return lines * size * 1.22 * spacing / 72


# ---------------------------------------------------------------- primitives
def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE,
         line_w=1.0):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, p in enumerate(paras):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.alignment = p.get("align", align)
        par.line_spacing = p.get("spacing", 1.15)
        par.space_after = Pt(p.get("space_after", 0))
        run = par.add_run()
        run.text = p["t"]
        f = run.font
        f.name = FONT
        f.size = Pt(p["size"])
        f.bold = p.get("bold", False)
        f.italic = p.get("italic", False)
        f.color.rgb = p.get("color", INK)
    return tb


def block_h(paras, w) -> float:
    return sum(est_h(p["t"], w, p["size"], p.get("spacing", 1.15),
                     p.get("bold", False)) + p.get("space_after", 0) / 72
               for p in paras)


SCALES: list[tuple[str, float]] = []


def img(slide, name, x, y, col_w, width) -> float:
    """Place at an explicit width, centred in the column, return the height."""
    path = FIG / name
    iw, ih = Image.open(path).size
    h = width * ih / iw
    slide.shapes.add_picture(str(path), Inches(x + (col_w - width) / 2),
                             Inches(y), Inches(width), Inches(h))
    SCALES.append((name, width / (iw / 300)))
    return h


def para(slide, x, y, w, paras) -> float:
    h = block_h(paras, w)
    text(slide, x, y, w, h + 0.4, paras)
    return h


def caption(slide, x, y, w, lead, body, size=19.5) -> float:
    return para(slide, x, y, w, [
        {"t": lead, "size": size, "bold": True, "color": INK, "space_after": 2},
        {"t": body, "size": size, "color": INK_2},
    ])


def sec_header(slide, x, w, y, num, title) -> float:
    """Uniform dark section bar -- the conventional poster idiom."""
    bh = 1.32
    rect(slide, x, y, w, bh, fill=INK_D)
    rect(slide, x, y, 0.16, bh, fill=TEAL)
    text(slide, x + 0.60, y, 1.30, bh,
         [{"t": num, "size": 34, "bold": True, "color": TEAL}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(slide, x + 1.72, y, w - 2.10, bh,
         [{"t": title, "size": 31, "bold": True, "color": WHITE}],
         anchor=MSO_ANCHOR.MIDDLE)
    return bh


def sub_header(slide, x, w, y, title) -> float:
    """Lighter heading used for continuations inside a section."""
    h = para(slide, x, y, w,
             [{"t": title, "size": 25, "bold": True, "color": INK}])
    rect(slide, x, y + h + 0.16, w, 0.045, fill=TEAL)
    return h + 0.16 + 0.045


# ---------------------------------------------------------------- canvas
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])
rect(slide, 0, 0, W, H, fill=PAPER)

# ---------------------------------------------------------------- header
rect(slide, 0, 0, W, HDR_H, fill=INK_D)
text(slide, MARGIN, 0.62, 33.5, 2.4, [
    {"t": "A Mechanism-Informed Gene Signature for\n"
          "Regenerative vs. Fibrotic Tendon Healing",
     "size": 53, "bold": True, "color": WHITE, "spacing": 1.0},
])
text(slide, MARGIN, 3.16, 34.0, 1.9, [
    {"t": "Mapping the transcription-factor program behind the fate decision "
          "in tendon progenitors — and compressing it into a validated "
          "23-gene Healing Index.",
     "size": 20, "color": MUTED, "spacing": 1.12, "space_after": 10},
    {"t": "Fangyi Li¹    ·    Wilder Scott²    ·    Yan Yan³",
     "size": 22, "bold": True, "color": WHITE, "space_after": 4},
    {"t": "¹Department of Integrative Biology, College of Biological Science, "
          "University of Guelph    ·    ²Sunnybrook Research Institute    ·    "
          "³School of Computer Science, College of Computational, Mathematical "
          "and Physical Sciences, University of Guelph",
     "size": 16.0, "color": MUTED, "spacing": 1.14},
])
LOGOS = [
    (LOGO / "sunnybrook3.png", 7.60),
    (LOGO / "GUELPH/SINGLECOLOUR_DARKBG/"
            "CBS_DIB_GUELPH_SINGLECOLOUR_DARKBG_V1_1.5IN_H_300PPI.png", 6.05),
]
_gap = 0.90
_lx = W - MARGIN - sum(w for _, w in LOGOS) - _gap * (len(LOGOS) - 1)
for _path, _lw in LOGOS:
    _iw, _ih = Image.open(_path).size
    _lh = _lw * _ih / _iw
    slide.shapes.add_picture(str(_path), Inches(_lx), Inches((HDR_H - _lh) / 2),
                             Inches(_lw), Inches(_lh))
    _lx += _lw + _gap

# Results spans columns 3-4; tint the pair so the span reads as one section.
rect(slide, COL_X[2] - 0.34, BODY_Y - 0.30,
     2 * COL_W + GUT + 0.68, BODY_H + 0.56, fill=TINT)

# ---------------------------------------------------------------- 1. INTRO
x, cw = COL_X[0], COL_W
y = BODY_Y
y += sec_header(slide, x, cw, y, "1", "Introduction") + 0.34
y += para(slide, x, y, cw, [
    {"t": "Tendon heals badly: repair tissue is scar, not aligned tendon. "
          "Two progenitor populations inside the same PDGFRα⁺ pool decide "
          "which one forms — tendon stem/progenitor cells (TSPC) rebuild "
          "tendon, tendon fibro-adipogenic progenitors (T-FAP) lay down scar "
          "(Harvey et al., 2019).",
     "size": 20, "color": INK_2, "spacing": 1.16}]) + 0.30
y += img(slide, "fig1_umap_subclusters.png", x, y, cw, cw) + 0.20
y += caption(slide, x, y, cw,
             "3,900 PDGFRα⁺ cells from injured mouse tendon.",
             "Sub-clustering resolves four populations "
             "(Cherief et al., 2023; day 14 post-injury).") + 0.34
y += img(slide, "fig2_composition_shift.png", x, y, cw, cw) + 0.20
y += caption(slide, x, y, cw,
             "Blocking sensory-nerve signalling shifts the pool.",
             "TSPC share falls by more than half; T-FAP becomes the largest "
             "population. Nerve signalling is blocked chemical-genetically "
             "(TrkA-F592A mice + 1NMPP1), not surgically — “denervated” is "
             "shorthand throughout.") + 0.40
ob = [
    {"t": "OBJECTIVE", "size": 16, "bold": True, "color": TEAL,
     "space_after": 7},
    {"t": "Identify the transcription-factor program that carries out this "
          "innervation-dependent fate decision, then compress it into a "
          "minimal, validated gene signature that scores how regeneratively "
          "a tendon is healing.",
     "size": 21, "bold": True, "color": INK, "spacing": 1.14},
]
ob_h = block_h(ob, cw - 1.20) + 0.80
rect(slide, x, y, cw, ob_h, fill=BAND)
rect(slide, x, y, 0.14, ob_h, fill=TEAL)
text(slide, x + 0.62, y + 0.40, cw - 1.20, ob_h - 0.6, ob)
end1 = y + ob_h

# ---------------------------------------------------------------- 2. METHODS
x, cw = COL_X[1], COL_W
y = BODY_Y
y += sec_header(slide, x, cw, y, "2", "Methods") + 0.34
y += para(slide, x, y, cw, [
    {"t": "Two independent public mouse scRNA-seq datasets, analysed through "
          "one integrated pipeline. No new animal work was performed.",
     "size": 20, "color": INK_2, "spacing": 1.16}]) + 0.30

for name, meta, role in [
    ("Harvey et al. (2019)",
     "Patellar tendon, uninjured · 4,069 cells after QC · SRA PRJNA506218",
     "Ground-truth TSPC / T-FAP labels. Training set."),
    ("Cherief et al. (2023)",
     "Achilles tendon, day 14 post-injury · 22,615 cells · GEO GSE244921",
     "Innervated vs. TrkA-blocked. Held-out test set."),
]:
    card = [
        {"t": name, "size": 21, "bold": True, "color": INK, "space_after": 4},
        {"t": meta, "size": 16.5, "color": INK_2, "spacing": 1.12,
         "space_after": 5},
        {"t": role, "size": 17.5, "bold": True, "color": TEAL, "spacing": 1.12},
    ]
    ch = block_h(card, cw - 1.10) + 0.62
    rect(slide, x, y, cw, ch, fill=WHITE, line=LINE)
    rect(slide, x, y, 0.11, ch, fill=INK)
    text(slide, x + 0.52, y + 0.31, cw - 1.10, ch - 0.5, card)
    y += ch + 0.40
y += 0.28

y += para(slide, x, y, cw, [
    {"t": "Analysis pipeline", "size": 22, "bold": True, "color": INK}]) + 0.26
STEPS = [
    ("QC, clustering, annotation",
     "PDGFRα⁺ pool (3,900 cells) sub-clustered into four populations"),
    ("pySCENIC — TF regulon inference",
     "Run independently per dataset, then compared (Van de Sande et al., 2020)"),
    ("CellRank 2 — fate drivers",
     "Genes correlated with commitment to each terminal population "
     "(Weiler et al., 2024)"),
    ("LIANA — cell–cell communication",
     "Ligand–receptor signalling, innervated vs. blocked "
     "(Dimitrov et al., 2022)"),
    ("Feature selection",
     "LASSO, Boruta and XGBoost run separately on 124 mechanism-derived "
     "candidates"),
    ("Classifier and Healing Index",
     "Logistic regression trained on Harvey, applied unchanged to Cherief"),
]
for i, (head, body) in enumerate(STEPS, 1):
    rect(slide, x + 0.02, y + 0.06, 0.62, 0.62, fill=TEAL,
         shape=MSO_SHAPE.OVAL)
    text(slide, x + 0.02, y + 0.06, 0.62, 0.62,
         [{"t": str(i), "size": 17, "bold": True, "color": WHITE}],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    sp = [
        {"t": head, "size": 19, "bold": True, "color": INK, "space_after": 2},
        {"t": body, "size": 16.5, "color": INK_2, "spacing": 1.12},
    ]
    sh = block_h(sp, cw - 0.90)
    text(slide, x + 0.90, y, cw - 0.90, sh + 0.4, sp)
    y += max(sh, 0.70) + 0.44
y += 0.26

y += para(slide, x, y, cw, [
    {"t": "Robustness checks", "size": 22, "bold": True, "color": INK,
     "space_after": 4},
    {"t": "Every headline claim was re-tested. Two did not survive.",
     "size": 17, "color": INK_2, "spacing": 1.12}]) + 0.26
CHECKS = [
    ('"Efna1→Epha3 is the strongest signal"',
     "LIANA permutation test, n = 100", "RETRACTED", "p = 1.0", RED),
    ('"T-FAP regulons feed back on Pdgfra"',
     "Hypergeometric, BH-corrected", "RETRACTED", "nothing at p < 0.05", RED),
    ('"Zero TSPC regulons replicate"',
     "200× power-matched stability selection", "CONFIRMED",
     "still zero at matched n", TEAL),
]
for claim, test, verdict, detail, col in CHECKS:
    cp = [
        {"t": claim, "size": 17.5, "bold": True, "color": INK, "spacing": 1.08,
         "space_after": 3},
        {"t": test, "size": 15.5, "italic": True, "color": INK_2,
         "space_after": 4},
        {"t": f"{verdict}  ·  {detail}", "size": 17, "bold": True,
         "color": col},
    ]
    ch = block_h(cp, cw - 0.80) + 0.44
    rect(slide, x, y, 0.09, ch, fill=col)
    text(slide, x + 0.42, y + 0.22, cw - 0.80, ch - 0.4, cp)
    y += ch + 0.32
end2 = y

# ---------------------------------------------------------------- 3. RESULTS
x, cw = COL_X[2], COL_W
y = BODY_Y
y += sec_header(slide, x, cw, y, "3", "Results") + 0.34
y += para(slide, x, y, cw, [
    {"t": "Which transcription-factor regulons separate the two fates, and do "
          "they replicate in a second, independent dataset?",
     "size": 20, "color": INK_2, "spacing": 1.16}]) + 0.30
y += img(slide, "fig3_regulon_scatter.png", x, y, cw, cw) + 0.20
y += caption(slide, x, y, cw,
             "Fifteen regulons replicate as fibrotic. None replicate as "
             "regenerative.",
             "A 200-iteration power-matched audit reproduced the empty "
             "quadrant, so it is not a sample-size artefact.") + 0.34
y += img(slide, "fig4_circuit.png", x, y, cw, cw) + 0.18
y += caption(slide, x, y, cw,
             "The surviving program is an AP-1 / KLF / NF-κB network.",
             "Junb, Jund and Klf9 contribute the most fate-driver genes; 14 of "
             "the 23 final panel genes are their targets. 12 of the 15 "
             "regulons are shown — three contribute none.") + 0.40
kf = [
    {"t": "KEY FINDING", "size": 16, "bold": True, "color": SALMON,
     "space_after": 7},
    {"t": "The fate decision is asymmetric. Innervation appears to suppress a "
          "default fibrotic program rather than to activate a regenerative "
          "one — the mirror image of our starting hypothesis.",
     "size": 21, "bold": True, "color": WHITE, "spacing": 1.14},
]
kf_h = block_h(kf, cw - 1.20) + 0.80
rect(slide, x, y, cw, kf_h, fill=INK_D)
rect(slide, x, y, 0.14, kf_h, fill=RED)
text(slide, x + 0.62, y + 0.40, cw - 1.20, kf_h - 0.6, kf)
end3 = y + kf_h

# ------------------------------------------------- 3. RESULTS (continued)
x, cw = COL_X[3], COL_W
y = BODY_Y + 1.32 + 0.34          # align with column 3's post-header content
y += sub_header(slide, x, cw, y, "From mechanism to a gene signature") + 0.30
y += img(slide, "fig5_venn.png", x, y, cw, 5.95) + 0.22
y += caption(slide, x, y, cw,
             "Three methods agree on 23 genes.",
             "LASSO, Boruta and XGBoost ran independently on 124 "
             "mechanism-derived candidates.") + 0.30
y += img(slide, "fig6_roc.png", x, y, cw, 5.65) + 0.22
y += caption(slide, x, y, cw,
             "The panel generalises to a held-out dataset.",
             "Trained on Harvey et al. (2019), then applied unchanged to "
             "Cherief et al. (2023) — independently labelled, never seen "
             "during selection or training.") + 0.30
y += img(slide, "fig7_healing_index.png", x, y, cw, cw) + 0.20
y += caption(slide, x, y, cw,
             "The Healing Index tracks innervation status.",
             "It also places two cell identities it never saw in training "
             "(Stromal, Tenogenic-progenitor) between the two fates.")
end4 = y

# ---------------------------------------------------------------- 4. CONCLUSIONS
CY = BODY_Y + BODY_H + 0.46
CH = 3.30
rect(slide, 0, CY, W, CH, fill=BAND)
ch_h = sec_header(slide, MARGIN, 13.2, CY + 0.34, "4", "Conclusions")

cx = MARGIN + 13.9
text(slide, MARGIN, CY + 0.34 + ch_h + 0.26, 13.2, 1.9, [
    {"t": "A mechanism-derived 23-gene Healing Index now tracks the "
          "regenerative-to-fibrotic shift across two independent tendon "
          "datasets.",
     "size": 19.5, "bold": True, "color": INK, "spacing": 1.14},
])
text(slide, cx, CY + 0.46, 17.6, CH - 0.9, [
    {"t": "▪  Innervation suppresses a fibrotic program rather than activating "
          "a regenerative one: 15 T-FAP regulons replicate across datasets, "
          "0 TSPC regulons do.", "size": 18, "color": INK, "spacing": 1.12,
     "space_after": 7},
    {"t": "▪  A 23-gene panel separates the two fates at 0.989 ROC-AUC on a "
          "held-out, independently labelled dataset.", "size": 18,
     "color": INK, "spacing": 1.12, "space_after": 7},
    {"t": "▪  The Healing Index falls from 25.2 to 15.5 when nerve signalling "
          "is blocked, and orders unseen cell identities sensibly.",
     "size": 18, "color": INK, "spacing": 1.12},
])
lx = cx + 18.3
text(slide, lx, CY + 0.46, W - MARGIN - lx, CH - 0.9, [
    {"t": "LIMITATIONS  ·  WHAT'S NEXT", "size": 15.5, "bold": True,
     "color": INK_2, "space_after": 6},
    {"t": "▪  No reproducible regenerative regulon is not proof none exists — "
          "both datasets are small.", "size": 17, "color": INK, "spacing": 1.1,
     "space_after": 5},
    {"t": "▪  Neither dataset has biological replicates — validation is "
          "cross-dataset, not cross-animal.", "size": 17, "color": INK,
     "spacing": 1.1, "space_after": 5},
    {"t": "▪  TSPC recall is 77–78% vs. 98–99% for T-FAP: the index flags "
          "fibrotic drift better than it confirms repair.", "size": 17,
     "color": INK, "spacing": 1.1, "space_after": 5},
    {"t": "▪  Next: score newly generated lab tendon scRNA-seq — the first "
          "true biological-replicate test.", "size": 17, "color": INK,
     "spacing": 1.1},
])

# ---------------------------------------------------------------- footer
FY = CY + CH + 0.40
rect(slide, MARGIN, FY, W - 2 * MARGIN, 0.035, fill=LINE)
fy = FY + 0.24
for fx, fw, head, body in [
    (MARGIN, 21.6, "Data & code availability",
     "All data are publicly available (GEO GSE244921; SRA PRJNA506218). No "
     "unpublished or confidential data were used. Computation was performed on "
     "the Digital Research Alliance of Canada (NIBI). Code and analysis "
     "notebooks: github.com/lifangy6/sunnybrook-tendon-fate-mapping"),
    (MARGIN + 22.8, W - 2 * MARGIN - 22.8, "Acknowledgements",
     "We thank Dr. Karl Cottenie for course instruction and for feedback at "
     "every stage of this project, and Sabrina Saiphoo for her generosity with "
     "help and shared resources during our time in the same lab."),
]:
    text(slide, fx, fy, fw, 1.3, [
        {"t": head, "size": 16, "bold": True, "color": INK, "space_after": 3},
        {"t": body, "size": 15, "color": INK_2, "spacing": 1.15},
    ])

fy += 1.18
REFS = [
    "Cherief, M., et al. (2023). TrkA-mediated sensory innervation of injured "
    "mouse tendon supports tendon sheath progenitor cell expansion and tendon "
    "repair. Science Translational Medicine, 15(727), eade4619.",
    "Dimitrov, D., et al. (2022). Comparison of methods and resources for "
    "cell-cell communication inference from single-cell RNA-Seq data. Nature "
    "Communications, 13, 3224.",
    "Harvey, T., Flamenco, S., & Fan, C.-M. (2019). A Tppp3+Pdgfrα+ tendon stem "
    "cell population contributes to regeneration and reveals a shared role for "
    "PDGF signalling in regeneration and fibrosis. Nature Cell Biology, 21, "
    "1490–1503.",
    "Van de Sande, B., et al. (2020). A scalable SCENIC workflow for "
    "single-cell gene regulatory network analysis. Nature Protocols, 15, "
    "2247–2276.",
    "Weiler, P., Lange, M., Klein, M., Pe'er, D., & Theis, F. J. (2024). "
    "CellRank 2: Unified fate mapping in multiview single-cell data. Nature "
    "Methods, 21, 1196–1205.",
]
text(slide, MARGIN, fy, 5.4, 0.6,
     [{"t": "References", "size": 16, "bold": True, "color": INK}])
ref_w = (W - 2 * MARGIN - 5.6 - 2 * 0.7) / 3
for i, r in enumerate(REFS):
    cxr = MARGIN + 5.6 + (i % 3) * (ref_w + 0.7)
    cyr = fy + (i // 3) * 0.62
    text(slide, cxr, cyr, ref_w, 0.60,
         [{"t": f"{i + 1}.  {r}", "size": 13, "color": INK_2, "spacing": 1.08}])


OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)

limit = BODY_Y + BODY_H
print(f"saved {OUT}")
print(f"\ncolumn budget (must stay under {limit:.2f} in):")
for lbl, e in [("1 Introduction", end1), ("2 Methods", end2),
               ("3 Results", end3), ("3 Results cont.", end4)]:
    print(f"  {lbl:17s} ends {e:6.2f}   headroom {limit - e:+.2f}"
          f"   {'OK' if e <= limit else '<-- OVERFLOW'}")
print(f"\nconclusions {CY:.2f}–{CY + CH:.2f} | footer {FY:.2f}–{fy + 1.10:.2f} "
      f"of {H}")
print("\nfigure scale vs. natural size (1.00 keeps point sizes literal):")
for n, s in SCALES:
    flag = "" if 0.85 <= s <= 1.35 else "  <-- check legibility"
    print(f"  {n:34s} {s:.2f}{flag}")
