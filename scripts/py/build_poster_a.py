"""Assemble PLAN A: the three-act narrative 48 x 36 in poster.

A three-act argument reading left to right, a full-width band reporting the
findings that were stress-tested (two retracted, one confirmed), and a closing
walk-away statement. The conventional four-section alternative is
build_poster_b.py; the brainstorming the Poster Draft rubric asks for lives in
BINF6999/drafts/poster_brainstorm.md.

python-pptx cannot measure rendered text, so text-block heights are estimated
from character counts and the running y-position of each column is tracked
explicitly. The script prints a per-column budget at the end; if a column
exceeds ACT_H it will collide with the stress band, so check that output.

Run:  uv run --python 3.12 --with python-pptx --with pillow python \
          scripts/py/build_poster_a.py
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
FIG = ROOT / "figures" / "poster"
LOGO = ROOT / "BINF6999" / "drafts" / "logos"
OUT = ROOT / "BINF6999" / "drafts" / "LiF_BINF6999_Poster_Draft_A.pptx"

W, H = 48.0, 36.0
FONT = "Segoe UI"

INK = RGBColor(0x16, 0x20, 0x2A)
INK_D = RGBColor(0x0E, 0x16, 0x1E)
INK_2 = RGBColor(0x5C, 0x6B, 0x7A)
BLUE = RGBColor(0x1F, 0x6F, 0xA8)
RED = RGBColor(0xBE, 0x3A, 0x34)
SALMON = RGBColor(0xE8, 0x8A, 0x84)
PAPER = RGBColor(0xFC, 0xFC, 0xFA)
BAND = RGBColor(0xEF, 0xF2, 0xF1)
LINE = RGBColor(0xD8, 0xDE, 0xE2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x9A, 0xAA, 0xB8)
TINT2 = RGBColor(0xF8, 0xF2, 0xF1)

MARGIN, GUT = 1.30, 0.70
COL_WS = [14.0, 16.0, 14.0]
COL_X = [MARGIN,
         MARGIN + COL_WS[0] + GUT,
         MARGIN + COL_WS[0] + COL_WS[1] + 2 * GUT]

HDR_H = 4.85
ACT_Y = 5.35
ACT_H = 21.30
GAP = 0.24


# ---------------------------------------------------------------- text metrics
def est_h(txt: str, w: float, size: float, spacing: float = 1.15,
          bold: bool = False) -> float:
    char_w = 0.545 if bold else 0.495
    cpl = max(1, int(w * 72 / (char_w * size)))
    lines = sum(max(1, -(-len(seg) // cpl)) for seg in txt.split("\n"))
    return lines * size * 1.22 * spacing / 72


# ---------------------------------------------------------------- primitives
def rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE,
         line_w=1.0):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, p in enumerate(paras):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.alignment = p.get("align", align)
        par.line_spacing = p.get("spacing", 1.15)
        par.space_after = Pt(p.get("space_after", 0))
        run = par.add_run()
        run.text = p["t"]
        f = run.font
        f.name = FONT
        f.size = Pt(p["size"])
        f.bold = p.get("bold", False)
        f.italic = p.get("italic", False)
        f.color.rgb = p.get("color", INK)
    return tb


def block_h(paras, w) -> float:
    return sum(est_h(p["t"], w, p["size"], p.get("spacing", 1.15),
                     p.get("bold", False)) + p.get("space_after", 0) / 72
               for p in paras)


def img(slide, name, x, y, col_w, width) -> float:
    """Place at an explicit width, centred in the column, and return the height.

    Figures are rendered at 300 dpi with bbox_inches='tight', so scale relative
    to their natural size is width / (px / 300). Keeping that near 1 is what
    keeps the point sizes chosen in poster_figures.py literal on the poster.
    """
    path = FIG / name
    iw, ih = Image.open(path).size
    h = width * ih / iw
    slide.shapes.add_picture(str(path), Inches(x + (col_w - width) / 2),
                             Inches(y), Inches(width), Inches(h))
    SCALES.append((name, width / (iw / 300)))
    return h


SCALES: list[tuple[str, float]] = []


def caption(slide, x, y, w, lead, body, size=20.5) -> float:
    paras = [
        {"t": lead, "size": size, "bold": True, "color": INK, "space_after": 2},
        {"t": body, "size": size, "color": INK_2},
    ]
    h = block_h(paras, w)
    text(slide, x, y, w, h + 0.4, paras)
    return h


def act_header(slide, i, n, title, accent, lead) -> float:
    x, cw = COL_X[i], COL_WS[i]
    d = 1.40
    rect(slide, x, ACT_Y, d, d, fill=accent, shape=MSO_SHAPE.OVAL)
    text(slide, x, ACT_Y + 0.19, d, d,
         [{"t": n, "size": 42, "bold": True, "color": WHITE}],
         align=PP_ALIGN.CENTER)
    t_par = [{"t": title, "size": 41, "bold": True, "color": INK, "spacing": 1.0}]
    th = block_h(t_par, cw - 1.78)
    text(slide, x + 1.78, ACT_Y - 0.12, cw - 1.78, th + 0.4, t_par)
    ry = ACT_Y + max(d + 0.18, th + 0.18)
    rect(slide, x, ry, cw, 0.05, fill=accent)
    l_par = [{"t": lead, "size": 23.0, "color": INK_2, "spacing": 1.16}]
    lh = block_h(l_par, cw)
    text(slide, x, ry + 0.32, cw, lh + 0.4, l_par)
    return (ry + 0.32 + lh) - ACT_Y


prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])
rect(slide, 0, 0, W, H, fill=PAPER)

rect(slide, 0, 0, W, HDR_H, fill=INK_D)
text(slide, MARGIN, 0.48, 33.0, 2.6, [
    {"t": "Tendon heals to scar by default.", "size": 64, "bold": True,
     "color": WHITE, "spacing": 0.96},
    {"t": "Nerves are what hold it back.", "size": 64, "bold": True,
     "color": SALMON, "spacing": 0.96},
])
text(slide, MARGIN, 3.14, 34.0, 1.9, [
    {"t": "Mapping the transcription-factor program behind the "
          "regenerative-vs-fibrotic fate decision in tendon progenitors — "
          "and compressing it into a validated 23-gene Healing Index.",
     "size": 23, "color": MUTED, "spacing": 1.12, "space_after": 11},
    {"t": "Fangyi Li¹    ·    Wilder Scott²    ·    Yan Yan³",
     "size": 23, "bold": True, "color": WHITE, "space_after": 4},
    {"t": "¹Department of Integrative Biology, College of Biological Science, "
          "University of Guelph    ·    ²Sunnybrook Research Institute    ·    "
          "³School of Computer Science, College of Computational, Mathematical "
          "and Physical Sciences, University of Guelph",
     "size": 16.5, "color": MUTED, "spacing": 1.14},
])
# Both logos, right-aligned and each vertically centred in the header band.
# The header text block runs out around x = 29, so the run can start at 32.
# Both are the pure-white single-colour variants (mean luminance 255), so the
# two marks match each other against the near-black header. The Guelph
# FULLCOLOUR_BLACKBG file is *not* interchangeable -- it keeps the crest's
# colour accents and reads as slightly off-white next to Sunnybrook.
LOGOS = [
    (LOGO / "sunnybrook3.png", 7.60),
    (LOGO / "GUELPH/SINGLECOLOUR_DARKBG/"
            "CBS_DIB_GUELPH_SINGLECOLOUR_DARKBG_V1_1.5IN_H_300PPI.png", 6.05),
]
_gap = 0.90
_lx = W - MARGIN - sum(w for _, w in LOGOS) - _gap * (len(LOGOS) - 1)
for _path, _lw in LOGOS:
    _iw, _ih = Image.open(_path).size
    _lh = _lw * _ih / _iw
    slide.shapes.add_picture(str(_path), Inches(_lx), Inches((HDR_H - _lh) / 2),
                             Inches(_lw), Inches(_lh))
    _lx += _lw + _gap

rect(slide, COL_X[1] - 0.34, ACT_Y - 0.30, COL_WS[1] + 0.68, ACT_H + 0.52,
     fill=TINT2)

# ---------------------------------------------------------------- act 1
x, cw = COL_X[0], COL_WS[0]
y = ACT_Y + act_header(
    slide, 0, "1", "Two cells, one decision", INK,
    lead="Tendon stem/progenitor cells (TSPC) rebuild aligned tendon. "
         "Tendon fibro-adipogenic progenitors (T-FAP) lay down scar instead. "
         "Both live in the same PDGFRα⁺ pool.")
y += GAP + 0.10
y += img(slide, "fig1_umap_subclusters.png", x, y, cw, 11.65) + 0.20
y += caption(slide, x, y, cw,
             "3,900 PDGFRα⁺ cells from injured mouse tendon.",
             "Sub-clustering resolves four populations (Cherief et al., 2023; "
             "day 14 post-injury).") + 0.32
y += img(slide, "fig2_composition_shift.png", x, y, cw, 12.0) + 0.20
y += caption(slide, x, y, cw,
             "Blocking sensory-nerve signalling shifts the pool.",
             "TSPC share falls by more than half; T-FAP becomes the largest "
             "population. Nerve signalling is blocked chemical-genetically "
             "(TrkA-F592A mice + 1NMPP1), not surgically — “denervated” is "
             "shorthand throughout.") + 0.36
rect(slide, x, y, cw, 0.045, fill=LINE)
close1 = [{"t": "Innervation controls the outcome. The gene program that "
                "carries it out had never been mapped.",
           "size": 24, "bold": True, "color": INK, "spacing": 1.14}]
text(slide, x, y + 0.32, cw, block_h(close1, cw) + 0.4, close1)
end1 = y + 0.32 + block_h(close1, cw)

# ---------------------------------------------------------------- act 2
x, cw = COL_X[1], COL_WS[1]
y = ACT_Y + act_header(
    slide, 1, "2", "We looked for a repair switch.\nThere wasn't one.", RED,
    lead="Which transcription-factor regulons — a TF plus the genes it "
         "directly regulates — separate the two fates? Networks were inferred "
         "independently in each dataset with pySCENIC (Van de Sande et al., "
         "2020), then compared.")
y += GAP + 0.10
y += img(slide, "fig3_regulon_scatter.png", x, y, cw, 10.20) + 0.20
y += caption(slide, x, y, cw,
             "Fifteen regulons replicate as fibrotic. None replicate as "
             "regenerative.",
             "A 200-iteration power-matched audit reproduced the empty "
             "quadrant, so it is not a sample-size artefact.") + 0.32
y += img(slide, "fig4_circuit.png", x, y, cw, 13.00) + 0.18
y += caption(slide, x, y, cw,
             "The surviving program is an AP-1 / KLF / NF-κB network.",
             "Junb, Jund and Klf9 contribute the most fate-driver genes "
             "(CellRank; Weiler et al., 2024); 14 of the 23 final panel genes "
             "are their targets. 12 of the 15 regulons are shown — three "
             "contribute none.") + 0.34

rf = [
    {"t": "Fibrosis is not what happens when repair fails.", "size": 26,
     "bold": True, "color": WHITE, "spacing": 1.08, "space_after": 5},
    {"t": "Fibrosis is the default — and innervation is what suppresses it.",
     "size": 26, "bold": True, "color": SALMON, "spacing": 1.08, "space_after": 8},
    {"t": "The mirror image of our starting hypothesis, which predicted an "
          "innervation-activated regenerative program.",
     "size": 17, "italic": True, "color": MUTED, "spacing": 1.1},
]
rf_h = block_h(rf, cw - 1.30) + 0.76
rect(slide, x, y, cw, rf_h, fill=INK_D)
rect(slide, x, y, 0.14, rf_h, fill=RED)
text(slide, x + 0.66, y + 0.42, cw - 1.30, rf_h - 0.6, rf)
end2 = y + rf_h

# ---------------------------------------------------------------- act 3
x, cw = COL_X[2], COL_WS[2]
y = ACT_Y + act_header(
    slide, 2, "3", "So we measured the default", BLUE,
    lead="The fibrotic program is the one with reproducible signal, so it is "
         "also the one worth turning into a score.")
y += GAP + 0.10
# Venn and ROC side by side, scaled to a common height.
row_h = 5.85
vw = row_h * 1913 / 1565
rw = row_h * 1789 / 1660
gap3 = cw - vw - rw
img(slide, "fig5_venn.png", x, y, vw, vw)
img(slide, "fig6_roc.png", x + vw + gap3, y + 0.05, rw, rw)
y += row_h + 0.18
y += caption(slide, x, y, cw,
             "Three methods agree on 23 genes; the panel then generalises.",
             "LASSO, Boruta and XGBoost ran independently on 124 "
             "mechanism-derived candidates. Trained on Harvey et al. (2019) "
             "and applied unchanged to Cherief et al. (2023), whose labels "
             "were assigned independently from marker expression and never "
             "seen during selection or training.") + 0.34
y += img(slide, "fig7_healing_index.png", x, y, cw, 14.00) + 0.18
y += caption(slide, x, y, cw,
             "The score tracks innervation status.",
             "It also places two cell identities it never saw in training "
             "(Stromal, Tenogenic-progenitor) between the two fates.")
end3 = y

# ---------------------------------------------------------------- stress band
SY = ACT_Y + ACT_H + 0.42
SH = 2.62
rect(slide, 0, SY, W, SH, fill=BAND)
text(slide, MARGIN, SY + 0.52, 9.2, 2.2, [
    {"t": "We stress-tested our own headlines.", "size": 28, "bold": True,
     "color": INK, "spacing": 1.05, "space_after": 6},
    {"t": "Two of the three did not survive — and saying so is what makes the "
          "third worth trusting.", "size": 18, "color": INK_2, "spacing": 1.12},
])
cards = [
    ('"Efna1→Epha3 is the strongest differential signal"',
     "LIANA permutation test, n = 100 (Dimitrov et al., 2022)",
     "RETRACTED", "p = 1.0", RED),
    ('"T-FAP regulons feed back on Pdgfra"',
     "Hypergeometric enrichment, BH-corrected", "RETRACTED",
     "nothing at p < 0.05", RED),
    ('"Zero TSPC regulons replicate across datasets"',
     "200× power-matched stability selection", "CONFIRMED",
     "still zero at matched n", BLUE),
]
cx0, cgap = 11.30, 0.55
cw = (W - MARGIN - cx0 - 2 * cgap) / 3
for i, (claim, test, verdict, detail, col) in enumerate(cards):
    cx = cx0 + i * (cw + cgap)
    rect(slide, cx, SY + 0.34, cw, SH - 0.68, fill=WHITE, line=LINE)
    rect(slide, cx, SY + 0.34, 0.11, SH - 0.68, fill=col)
    text(slide, cx + 0.50, SY + 0.46, cw - 0.88, SH - 0.92, [
        {"t": claim, "size": 19.5, "bold": True, "color": INK, "spacing": 1.08,
         "space_after": 6},
        {"t": test, "size": 16.5, "italic": True, "color": INK_2,
         "space_after": 8},
        {"t": f"{verdict}   ·   {detail}", "size": 19.5, "bold": True,
         "color": col},
    ], anchor=MSO_ANCHOR.MIDDLE)

# ---------------------------------------------------------------- walk-away
BY = SY + SH + 0.40
BH = 2.48
rect(slide, MARGIN, BY, 26.9, BH, fill=INK_D)
rect(slide, MARGIN, BY, 0.15, BH, fill=BLUE)
text(slide, MARGIN + 0.68, BY + 0.38, 25.6, BH - 0.5, [
    {"t": "WHAT TO TAKE AWAY", "size": 16, "bold": True, "color": MUTED,
     "space_after": 8},
    {"t": "Healing a tendon may be less about switching regeneration on than "
          "about keeping the fibrotic default switched off. The Healing Index "
          "turns that idea into a number — ready to run on new tendon data the "
          "day it exists.",
     "size": 26, "bold": True, "color": WHITE, "spacing": 1.14},
])
lx = MARGIN + 27.6
text(slide, lx, BY + 0.06, W - MARGIN - lx, BH, [
    {"t": "LIMITATIONS  ·  WHAT'S NEXT", "size": 16, "bold": True,
     "color": INK_2, "space_after": 7},
    {"t": "▪  Finding no reproducible regenerative regulon is not proof that "
          "none exists — both datasets are small, and a weaker TSPC program "
          "could be missed.", "size": 17.5, "color": INK, "spacing": 1.1,
     "space_after": 6},
    {"t": "▪  Neither dataset has biological replicates — every validation so "
          "far is cross-dataset, not cross-animal.", "size": 17.5,
     "color": INK, "spacing": 1.1, "space_after": 6},
    {"t": "▪  TSPC recall is 77–78% against 98–99% for T-FAP: the index flags "
          "fibrotic drift more reliably than it confirms full repair.",
     "size": 17.5, "color": INK, "spacing": 1.1, "space_after": 6},
    {"t": "▪  Next: score newly generated lab tendon scRNA-seq — the first "
          "true biological-replicate test.", "size": 17.5, "color": INK,
     "spacing": 1.1},
])

# ---------------------------------------------------------------- footer
FY = BY + BH + 0.38
rect(slide, MARGIN, FY, W - 2 * MARGIN, 0.035, fill=LINE)
fy = FY + 0.24
for fx, fw, head, body in [
    (MARGIN, 21.6, "Data & code availability",
     "All data are publicly available (GEO GSE244921; SRA PRJNA506218). No "
     "unpublished or confidential data were used. Computation was performed on "
     "the Digital Research Alliance of Canada (NIBI). Code and analysis "
     "notebooks: github.com/lifangy6/sunnybrook-tendon-fate-mapping"),
    (MARGIN + 22.8, W - 2 * MARGIN - 22.8, "Acknowledgements",
     "We thank Dr. Karl Cottenie for course instruction and for feedback at "
     "every stage of this project, and Sabrina Saiphoo for her generosity with "
     "help and shared resources during our time in the same lab."),
]:
    text(slide, fx, fy, fw, 1.3, [
        {"t": head, "size": 16, "bold": True, "color": INK, "space_after": 3},
        {"t": body, "size": 15, "color": INK_2, "spacing": 1.15},
    ])

fy += 1.18
REFS = [
    "Cherief, M., et al. (2023). TrkA-mediated sensory innervation of injured "
    "mouse tendon supports tendon sheath progenitor cell expansion and tendon "
    "repair. Science Translational Medicine, 15(727), eade4619.",
    "Dimitrov, D., et al. (2022). Comparison of methods and resources for "
    "cell-cell communication inference from single-cell RNA-Seq data. Nature "
    "Communications, 13, 3224.",
    "Harvey, T., Flamenco, S., & Fan, C.-M. (2019). A Tppp3+Pdgfrα+ tendon stem "
    "cell population contributes to regeneration and reveals a shared role for "
    "PDGF signalling in regeneration and fibrosis. Nature Cell Biology, 21, "
    "1490–1503.",
    "Van de Sande, B., et al. (2020). A scalable SCENIC workflow for "
    "single-cell gene regulatory network analysis. Nature Protocols, 15, "
    "2247–2276.",
    "Weiler, P., Lange, M., Klein, M., Pe'er, D., & Theis, F. J. (2024). "
    "CellRank 2: Unified fate mapping in multiview single-cell data. Nature "
    "Methods, 21, 1196–1205.",
]
text(slide, MARGIN, fy, 5.4, 0.6,
     [{"t": "References", "size": 16, "bold": True, "color": INK}])
ref_w = (W - 2 * MARGIN - 5.6 - 2 * 0.7) / 3
for i, r in enumerate(REFS):
    cxr = MARGIN + 5.6 + (i % 3) * (ref_w + 0.7)
    cyr = fy + (i // 3) * 0.62
    text(slide, cxr, cyr, ref_w, 0.60,
         [{"t": f"{i + 1}.  {r}", "size": 13, "color": INK_2, "spacing": 1.08}])


OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)

limit = ACT_Y + ACT_H
print(f"saved {OUT}")
print(f"\npage 1 column budget (must stay under {limit:.2f} in):")
for i, e in enumerate((end1, end2, end3), 1):
    print(f"  act {i}: ends {e:6.2f}   headroom {limit - e:+.2f}"
          f"   {'OK' if e <= limit else '<-- OVERFLOW'}")
print(f"\nstress {SY:.2f}–{SY + SH:.2f} | bottom {BY:.2f}–{BY + BH:.2f} "
      f"| footer {FY:.2f}–{fy + 1.10:.2f} of {H}")
print("\nfigure scale vs. natural size (1.00 keeps point sizes literal):")
for n, s in SCALES:
    flag = "" if 0.85 <= s <= 1.45 else "   <-- check legibility"
    print(f"  {n:32s} {s:5.2f}x{flag}")
